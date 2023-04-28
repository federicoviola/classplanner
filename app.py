from flask import Flask, render_template, request, send_file, redirect, url_for, Response
import openai
from pptx import Presentation
from dotenv import load_dotenv
import io
import os
import tempfile
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
from reportlab.lib.units import inch

openai_api_key = os.environ.get("OPENAI_API_KEY")
load_dotenv()
app = Flask(__name__)

@app.route('/')
def index():
    return render_template('index.html')

def generate_class_plan(topic, duration, education_level, user_keywords):
    keywords = ", ".join(user_keywords)  # user_keywords debe ser una lista de palabras clave proporcionadas por el usuario
    prompt = f"Crear una presentacion para una clase de nivel educativo {education_level}. que tenga como objetivo enseñar a los estudiantes sobre el tema '{topic}' con una duración de {duration} minutos. Incluir los siguientes conceptos clave: {keywords}."

    response = openai.Completion.create(
        engine="text-davinci-003",
        prompt=prompt,
        max_tokens=150,
        n=1,
        stop=None,
        temperature=0.5,
    )

    return response.choices[0].text.strip()

def generate_presentation(topic, plan_text):
    # Crea una nueva presentación
    ppt = Presentation()

    # Añade una diapositiva de título
    title_slide = ppt.slide_layouts[0]
    slide = ppt.slides.add_slide(title_slide)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "Clase: " + topic
    subtitle.text = "Plan de clase generado automáticamente"

    # Divide el plan de clase en puntos
    points = plan_text.split("\n")

    # Añade una diapositiva para cada punto
    for point in points:
        bullet_slide = ppt.slide_layouts[1]
        slide = ppt.slides.add_slide(bullet_slide)

        shapes = slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]

        title_shape.text = topic
        text_frame = body_shape.text_frame
        text_frame.text = point

    # Guarda la presentación en un objeto BytesIO
    presentation_data = io.BytesIO()
    ppt.save(presentation_data)
    presentation_data.seek(0)

    return presentation_data

def generate_pdf(topic, plan_text):
    # Crea un nuevo objeto BytesIO para guardar el PDF
    pdf_data = io.BytesIO()

    # Crea un nuevo archivo PDF
    c = canvas.Canvas(pdf_data, pagesize=letter)

    # Agrega el título de la clase
    c.setFont("Helvetica-Bold", 16)
    c.drawString(inch, 10 * inch, "Clase: " + topic)

    # Agrega el plan de clase
    c.setFont("Helvetica", 12)
    y = 9 * inch
    for point in plan_text.split("\n"):
        c.drawString(inch, y, point)
        y -= 0.5 * inch

    # Guarda y cierra el PDF
    c.save()

    # Rebobina el objeto BytesIO y devuelve los datos
    pdf_data.seek(0)
    return pdf_data

from flask import send_file


@app.route('/generate', methods=['POST'])
def generate():
    topic = request.form['topic']
    duration = request.form['duration']
    file_type = request.form['file_type']
    education_level = request.form.get('education_level')
    keywords_input = request.form['keywords']

    user_keywords = [keyword.strip() for keyword in keywords_input.split(',')]

    # Genera la presentacion para la clase
    class_plan = generate_class_plan(topic, duration, education_level,user_keywords)

    # Genera la presentación o el PDF, según corresponda
    if file_type == 'pptx':
        file_data = generate_presentation(topic, class_plan)
        file_ext = 'pptx'
    elif file_type == 'pdf':
        file_data = generate_pdf(topic, class_plan)
        file_ext = 'pdf'
    else:
        return "Tipo de archivo no válido", 400

    # Envía el archivo generado como archivo descargable
    return send_file(
    file_data,
    download_name=f"{topic}_presentation.{file_ext}",
    as_attachment=True,
    )

if __name__ == '__main__':
    app.run(debug=True)
